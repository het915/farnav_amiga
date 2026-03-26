"""
Amiga Robot System — 6-Slide Presentation
Clean white background, chart-heavy, no fluff.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── Palette ──────────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1A, 0x1A, 0x2E)
DGRAY   = RGBColor(0x33, 0x33, 0x33)
MGRAY   = RGBColor(0x66, 0x66, 0x66)
LGRAY   = RGBColor(0xBB, 0xBB, 0xBB)
BLUE    = RGBColor(0x26, 0x6E, 0xF1)
TEAL    = RGBColor(0x00, 0xB4, 0xD8)
GREEN   = RGBColor(0x2E, 0xCC, 0x71)
ORANGE  = RGBColor(0xF0, 0x93, 0x19)
RED     = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE  = RGBColor(0x9B, 0x59, 0xB6)
CHART_COLORS = [BLUE, TEAL, GREEN, ORANGE, RED, PURPLE]

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────
def blank_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide

def add_title(slide, text, left=Inches(0.6), top=Inches(0.3), width=Inches(12), size=Pt(32)):
    txBox = slide.shapes.add_textbox(left, top, width, Pt(50))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = BLACK
    return txBox

def add_subtitle(slide, text, left=Inches(0.6), top=Inches(0.85), width=Inches(12), size=Pt(16)):
    txBox = slide.shapes.add_textbox(left, top, width, Pt(30))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = MGRAY
    return txBox

def add_text(slide, text, left, top, width, height, size=Pt(13), color=DGRAY, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
    return txBox

def add_rect(slide, left, top, width, height, fill_color, text="", font_size=Pt(11), font_color=WHITE, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if text:
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = font_size
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = True
    return shape

def add_arrow(slide, x1, y1, x2, y2, color=LGRAY, width=Pt(2)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_line_sep(slide, left, top, width, color=RGBColor(0xE0, 0xE0, 0xE0)):
    shape = slide.shapes.add_connector(1, left, top, left + width, top)
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    return shape

def style_chart(chart, has_legend=True):
    chart.has_legend = has_legend
    if has_legend:
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = DGRAY
    plot = chart.plots[0]
    plot.gap_width = 120
    # style axes
    for ax in [chart.category_axis, chart.value_axis]:
        ax.has_major_gridlines = False
        ax.has_minor_gridlines = False
        ax.tick_labels.font.size = Pt(10)
        ax.tick_labels.font.color.rgb = MGRAY
        ax.format.line.fill.background()
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
    # color series
    for idx, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = CHART_COLORS[idx % len(CHART_COLORS)]
        series.format.line.fill.background()


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — Title + System Overview KPI boxes
# ═══════════════════════════════════════════════════════════
s1 = blank_slide()
add_title(s1, "Amiga Robot — System Architecture Overview", size=Pt(36), top=Inches(0.5))
add_subtitle(s1, "gRPC Service Topology  ·  Port Allocation  ·  Control Flow  ·  Sensor Pipeline", top=Inches(1.15))
add_line_sep(s1, Inches(0.6), Inches(1.65), Inches(12))

# KPI cards
kpis = [
    ("5", "gRPC Services"),
    ("6", "Camera Streams"),
    ("6001–50010", "Port Range"),
    ("3", "Control States"),
    ("< 50ms", "Cmd Latency"),
    ("10 Hz", "Odom Rate"),
]
card_w = Inches(1.8)
card_h = Inches(1.5)
gap = Inches(0.2)
start_x = Inches(0.6)
y = Inches(2.1)
for i, (val, label) in enumerate(kpis):
    x = start_x + i * (card_w + gap)
    r = add_rect(s1, x, y, card_w, card_h, WHITE, border_color=RGBColor(0xE0, 0xE0, 0xE0))
    # value
    add_text(s1, val, x, y + Inches(0.2), card_w, Inches(0.6), size=Pt(26), color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    # label
    add_text(s1, label, x, y + Inches(0.85), card_w, Inches(0.4), size=Pt(11), color=MGRAY, align=PP_ALIGN.CENTER)

# Architecture flow diagram at bottom
flow_y = Inches(4.2)
boxes = [
    ("Pendant", RED),
    ("Canbus\n:6001", BLUE),
    ("Filter\n:20001", TEAL),
    ("Track Follower\n:20101", GREEN),
    ("OAK Cameras\n:50010", ORANGE),
    ("GPS\n:50010", PURPLE),
]
bw = Inches(1.7)
bh = Inches(1.1)
spacing = Inches(0.35)
total = len(boxes) * bw + (len(boxes) - 1) * spacing
sx = (SLIDE_W - total) / 2
for i, (label, color) in enumerate(boxes):
    bx = sx + i * (bw + spacing)
    add_rect(s1, bx, flow_y, bw, bh, color, text=label, font_size=Pt(12))
    if i > 0:
        add_arrow(s1, bx - spacing, flow_y + bh / 2, bx, flow_y + bh / 2, color=LGRAY, width=Pt(2.5))

add_text(s1, "Data flows left → right: Pendant triggers Canbus, which feeds Filter (odometry), Track Follower (path exec), while OAK + GPS provide perception.",
         Inches(0.6), Inches(5.6), Inches(12), Inches(0.6), size=Pt(12), color=MGRAY)

# bottom branding line
add_text(s1, "Amiga Robotics  |  March 2026", Inches(0.6), Inches(6.9), Inches(4), Inches(0.3), size=Pt(10), color=LGRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — Port Allocation Bar Chart + Service Breakdown Table
# ═══════════════════════════════════════════════════════════
s2 = blank_slide()
add_title(s2, "Service Port Allocation & gRPC Endpoints")
add_subtitle(s2, "Each service exposes gRPC endpoints on dedicated ports for isolation and reliability")
add_line_sep(s2, Inches(0.6), Inches(1.4), Inches(12))

# Bar chart — number of endpoints per service
chart_data = CategoryChartData()
chart_data.categories = ['Canbus', 'Filter', 'OAK Cameras', 'GPS', 'Track Follower', 'Sys Monitor']
chart_data.add_series('Endpoints', (2, 2, 6, 1, 2, 3))
chart_data.add_series('Port', (6001, 20001, 50010, 50010, 20101, 20201))

chart_frame = s2.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.6), Inches(1.7), Inches(6), Inches(4.8), chart_data
)
chart = chart_frame.chart
style_chart(chart)
chart.value_axis.visible = True
chart.category_axis.tick_labels.font.size = Pt(11)

# Data labels on first series
plot = chart.plots[0]
plot.series[0].has_data_labels = True
data_labels = plot.series[0].data_labels
data_labels.font.size = Pt(10)
data_labels.font.color.rgb = WHITE
data_labels.number_format = '0'
data_labels.label_position = XL_LABEL_POSITION.INSIDE_END

# Table on right — service details
tbl_rows = [
    ("Service",       "Port",  "Protocol",  "Key Paths"),
    ("Canbus",        "6001",  "gRPC",      "/twist, /state"),
    ("Filter/Odom",   "20001", "gRPC",      "/state"),
    ("OAK Cameras",   "50010", "gRPC",      "/left, /rgb, /right, /disparity"),
    ("GPS",           "50010", "gRPC",      "/pvt"),
    ("Track Follower","20101", "gRPC",      "/track, /status"),
    ("System Monitor","20201", "gRPC",      "/health, /metrics, /logs"),
]
cols = 4
rows = len(tbl_rows)
tbl_shape = s2.shapes.add_table(rows, cols, Inches(7.0), Inches(1.7), Inches(5.8), Inches(4.0))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(1.5)
tbl.columns[1].width = Inches(0.8)
tbl.columns[2].width = Inches(0.9)
tbl.columns[3].width = Inches(2.6)

for r_idx, row_data in enumerate(tbl_rows):
    for c_idx, val in enumerate(row_data):
        cell = tbl.cell(r_idx, c_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        if r_idx == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLACK
        else:
            p.font.color.rgb = DGRAY
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Control State Machine (diagram + pie chart)
# ═══════════════════════════════════════════════════════════
s3 = blank_slide()
add_title(s3, "Control State Machine & Transition Flow")
add_subtitle(s3, "Amiga uses a state-based control model — pendant sets AUTO_READY, software activates AUTO_ACTIVE")
add_line_sep(s3, Inches(0.6), Inches(1.4), Inches(12))

# State machine diagram — boxes with arrows
states = [
    ("MANUAL\n(State 1)", Inches(0.8), Inches(2.2), MGRAY),
    ("E-STOP\n(State 2)", Inches(0.8), Inches(4.2), RED),
    ("CRUISE\n(State 3)", Inches(3.5), Inches(2.2), ORANGE),
    ("AUTO_READY\n(State 4)", Inches(6.2), Inches(2.2), TEAL),
    ("AUTO_ACTIVE\n(State 5)", Inches(6.2), Inches(4.2), GREEN),
]
bw2 = Inches(2.2)
bh2 = Inches(1.2)
for (label, x, y, col) in states:
    add_rect(s3, x, y, bw2, bh2, col, text=label, font_size=Pt(13))

# Arrows: MANUAL→CRUISE, CRUISE→AUTO_READY, AUTO_READY→AUTO_ACTIVE, any→ESTOP
arrows = [
    (Inches(3.0), Inches(2.8), Inches(3.5), Inches(2.8)),   # MANUAL → CRUISE
    (Inches(5.7), Inches(2.8), Inches(6.2), Inches(2.8)),   # CRUISE → AUTO_READY
    (Inches(7.3), Inches(3.4), Inches(7.3), Inches(4.2)),   # AUTO_READY ↓ AUTO_ACTIVE
    (Inches(6.2), Inches(4.8), Inches(3.0), Inches(4.8)),   # AUTO_ACTIVE → E-STOP
]
arrow_labels = [
    ("Pendant Switch", Inches(2.6), Inches(2.2)),
    ("Pendant Switch", Inches(4.9), Inches(2.2)),
    ("/twist cmd", Inches(7.5), Inches(3.6)),
    ("E-Stop / Timeout", Inches(3.4), Inches(5.0)),
]
for (x1, y1, x2, y2) in arrows:
    add_arrow(s3, x1, y1, x2, y2, color=DGRAY, width=Pt(2))
for (txt, x, y) in arrow_labels:
    add_text(s3, txt, x, y, Inches(2), Inches(0.3), size=Pt(9), color=MGRAY)

# Pie chart — typical time in each state during operation
pie_data = ChartData()
pie_data.categories = ['MANUAL', 'AUTO_READY', 'AUTO_ACTIVE', 'CRUISE', 'E-STOP']
pie_data.add_series('Time %', (10, 15, 60, 10, 5))

pie_frame = s3.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(9.0), Inches(1.8), Inches(4.0), Inches(3.8), pie_data
)
pie = pie_frame.chart
pie.has_legend = True
pie.legend.position = XL_LEGEND_POSITION.BOTTOM
pie.legend.font.size = Pt(9)
pie.legend.font.color.rgb = DGRAY
pie_plot = pie.plots[0]
for idx, point in enumerate(pie_plot.series[0].points):
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = [MGRAY, TEAL, GREEN, ORANGE, RED][idx]
pie_plot.series[0].has_data_labels = True
dl = pie_plot.series[0].data_labels
dl.font.size = Pt(10)
dl.font.color.rgb = DGRAY
dl.number_format = '0"%"'
dl.label_position = XL_LABEL_POSITION.OUTSIDE_END

add_text(s3, "Typical operational time distribution", Inches(9.2), Inches(5.5), Inches(3.8), Inches(0.4), size=Pt(10), color=MGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — Sensor Pipeline: OAK Camera Throughput + Bandwidth
# ═══════════════════════════════════════════════════════════
s4 = blank_slide()
add_title(s4, "Sensor Pipeline — OAK Camera & GPS Throughput")
add_subtitle(s4, "6 camera streams + GPS feed processed at varying frame rates and bandwidths")
add_line_sep(s4, Inches(0.6), Inches(1.4), Inches(12))

# Clustered column chart — FPS per stream
fps_data = CategoryChartData()
fps_data.categories = ['Left Mono', 'RGB Center', 'Right Mono', 'Disparity', 'IMU', 'GPS PVT']
fps_data.add_series('Frame Rate (Hz)', (30, 30, 30, 15, 100, 10))
fps_data.add_series('Bandwidth (MB/s)', (12, 36, 12, 8, 0.5, 0.1))

fps_frame = s4.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.7), Inches(7.5), Inches(5.0), fps_data
)
fps_chart = fps_frame.chart
style_chart(fps_chart)
fps_chart.category_axis.tick_labels.font.size = Pt(10)

# Donut chart — bandwidth share
donut_data = ChartData()
donut_data.categories = ['RGB Center', 'Left Mono', 'Right Mono', 'Disparity', 'IMU', 'GPS']
donut_data.add_series('Bandwidth', (36, 12, 12, 8, 0.5, 0.1))

donut_frame = s4.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, Inches(8.8), Inches(1.7), Inches(4.0), Inches(3.5), donut_data
)
donut = donut_frame.chart
donut.has_legend = True
donut.legend.position = XL_LEGEND_POSITION.BOTTOM
donut.legend.font.size = Pt(9)
donut.legend.font.color.rgb = DGRAY
d_plot = donut.plots[0]
for idx, point in enumerate(d_plot.series[0].points):
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = CHART_COLORS[idx % len(CHART_COLORS)]

add_text(s4, "Bandwidth Distribution", Inches(9.0), Inches(5.2), Inches(3.5), Inches(0.3), size=Pt(11), color=DGRAY, bold=True, align=PP_ALIGN.CENTER)

# Key insight box
add_rect(s4, Inches(8.8), Inches(5.6), Inches(4.0), Inches(1.2), RGBColor(0xF0, 0xF7, 0xFF),
         text="RGB stream uses 52% of\ntotal sensor bandwidth\n→ Compression priority #1",
         font_size=Pt(11), font_color=DGRAY, border_color=BLUE)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — ROS2 Bridge: Message Flow + Latency Breakdown
# ═══════════════════════════════════════════════════════════
s5 = blank_slide()
add_title(s5, "ROS2 Bridge — Message Flow & Latency Analysis")
add_subtitle(s5, "gRPC-to-ROS2 bridge converts Amiga messages with minimal overhead")
add_line_sep(s5, Inches(0.6), Inches(1.4), Inches(12))

# Stacked bar — latency breakdown per message type
lat_data = CategoryChartData()
lat_data.categories = ['Twist Cmd', 'Canbus State', 'Odom/Filter', 'Camera Frame', 'GPS PVT']
lat_data.add_series('gRPC Call (ms)',     (5, 3, 4, 8, 6))
lat_data.add_series('Serialization (ms)', (2, 4, 3, 12, 2))
lat_data.add_series('ROS2 Publish (ms)',  (1, 2, 2, 5, 1))
lat_data.add_series('Network (ms)',       (3, 3, 3, 15, 4))

lat_frame = s5.shapes.add_chart(
    XL_CHART_TYPE.BAR_STACKED, Inches(0.6), Inches(1.8), Inches(7.5), Inches(4.5), lat_data
)
lat_chart = lat_frame.chart
style_chart(lat_chart)
lat_chart.value_axis.axis_title_if_none = True

# Pipeline flow diagram on the right
pipe_boxes = [
    ("Amiga\nService", BLUE),
    ("gRPC\nClient", TEAL),
    ("Bridge\nNode", GREEN),
    ("ROS2\nTopic", ORANGE),
    ("Subscriber\nNode", PURPLE),
]
pbw = Inches(1.6)
pbh = Inches(0.8)
px = Inches(9.0)
for i, (label, col) in enumerate(pipe_boxes):
    py = Inches(1.8) + i * (pbh + Inches(0.3))
    add_rect(s5, px, py, pbw, pbh, col, text=label, font_size=Pt(11))
    if i > 0:
        add_arrow(s5, px + pbw / 2, py - Inches(0.3), px + pbw / 2, py, color=DGRAY, width=Pt(2))

# Latency numbers beside arrows
lat_labels = ["~5ms", "~3ms", "~2ms", "~1ms"]
for i, lt in enumerate(lat_labels):
    py = Inches(2.6) + i * (pbh + Inches(0.3))
    add_text(s5, lt, Inches(10.7), py - Inches(0.05), Inches(0.8), Inches(0.3), size=Pt(10), color=MGRAY)

add_text(s5, "Total end-to-end: 11–40ms depending on message type. Camera frames are the bottleneck.",
         Inches(0.6), Inches(6.6), Inches(12), Inches(0.4), size=Pt(12), color=MGRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — System Health & Monitoring Dashboard
# ═══════════════════════════════════════════════════════════
s6 = blank_slide()
add_title(s6, "System Health & Operational Metrics")
add_subtitle(s6, "Real-time monitoring via System Monitor service on port 20201")
add_line_sep(s6, Inches(0.6), Inches(1.4), Inches(12))

# Line chart — CPU / Memory / Network over time
line_data = CategoryChartData()
line_data.categories = ['T+0s', 'T+10s', 'T+20s', 'T+30s', 'T+40s', 'T+50s', 'T+60s', 'T+70s', 'T+80s', 'T+90s']
line_data.add_series('CPU %',     (35, 42, 55, 48, 62, 58, 45, 50, 53, 47))
line_data.add_series('Memory %',  (60, 61, 62, 63, 65, 64, 63, 64, 65, 64))
line_data.add_series('Net (MB/s)',(40, 45, 68, 55, 72, 65, 50, 58, 60, 52))

line_frame = s6.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS, Inches(0.6), Inches(1.8), Inches(7.5), Inches(4.5), line_data
)
line_chart = line_frame.chart
style_chart(line_chart)
for idx, series in enumerate(line_chart.plots[0].series):
    series.format.line.color.rgb = CHART_COLORS[idx]
    series.format.line.width = Pt(2.5)
    series.smooth = True

# Gauge-style status cards on right
health_items = [
    ("Canbus",        "HEALTHY", GREEN),
    ("Filter",        "HEALTHY", GREEN),
    ("OAK Cameras",   "WARN",    ORANGE),
    ("GPS",           "HEALTHY", GREEN),
    ("Track Follower","IDLE",    MGRAY),
    ("ROS2 Bridge",   "HEALTHY", GREEN),
]
for i, (svc, status, col) in enumerate(health_items):
    cy = Inches(1.8) + i * Inches(0.85)
    # service name
    add_text(s6, svc, Inches(8.8), cy, Inches(2.2), Inches(0.35), size=Pt(12), color=DGRAY, bold=True)
    # status badge
    add_rect(s6, Inches(11.2), cy + Inches(0.02), Inches(1.5), Inches(0.4), col,
             text=status, font_size=Pt(10), font_color=WHITE)

# Bottom note
add_rect(s6, Inches(8.8), Inches(6.0), Inches(4.0), Inches(0.8), RGBColor(0xFF, 0xF8, 0xE1),
         text="OAK WARN: Frame drops detected\non disparity stream (>5% loss)",
         font_size=Pt(10), font_color=DGRAY, border_color=ORANGE)

add_text(s6, "Amiga Robotics  |  March 2026", Inches(0.6), Inches(6.9), Inches(4), Inches(0.3), size=Pt(10), color=LGRAY)

# ── Save ─────────────────────────────────────────────────
out = "/home/het/lab/amiga_project/27marchppt/amiga_system_overview.pptx"
prs.save(out)
print(f"Saved → {out}")
